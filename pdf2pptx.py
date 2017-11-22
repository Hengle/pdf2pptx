from pptx import Presentation
from pptx.util import Inches
import os
import shutil
import argparse

parser = argparse.ArgumentParser(description="description goes here")
parser.add_argument('fname', type=str, help='input file name')
parser.add_argument('-o', type=str, help='output file name', required=True)
parser.add_argument('-r', type=str, help='resolution. Default is 300.',
                    required=False, default=300)
parser.add_argument('--width', type=str,
                    help='Width of slide.Default is 10.',
                    required=False, default=10)
parser.add_argument('--height', type=str,
                    help='Height of slide.Default is 7.5',
                    required=False, default=7.5)

args = parser.parse_args()


def main():
    fpath = args.fname
    out_file = args.o
    density = str(args.r)
    tmp_dir = 'tmp'
    width = args.width
    height = args.height
    while(True):
        if os.path.exists('./' + tmp_dir):
            tmp_dir = tmp_dir + '_'
        else:
            break
    os.mkdir('./' + tmp_dir)
    tmp_img = './' + tmp_dir + '/tmp.png'
    os.system('convert -density ' + density + ' ' + fpath + ' ' + tmp_img)
    pics_len = len(os.listdir('./' + tmp_dir))
    pics = ['./tmp-' + str(n) + '.png' for n in range(pics_len)]
    prs = Presentation()
    for pic in pics:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Inches(0)
        width = Inches(10)
        height = Inches(7.5)
        left = Inches(0)
        pic = slide.shapes.add_picture(
            './' + tmp_dir + '/' + pic, left, top, width=width, height=height)
    prs.save(out_file)
    shutil.rmtree('./' + tmp_dir)


if __name__ == "__main__":
    main()
